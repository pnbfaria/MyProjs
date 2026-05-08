import sys
import os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

prs = Presentation(os.path.join(r'C:\Users\pedro\OneDrive\Git\MyProjs', 'template.pptx'))

# Create dir for extracted images
img_dir = os.path.join(r'C:\Users\pedro\OneDrive\Git\MyProjs', 'pptx_images')
os.makedirs(img_dir, exist_ok=True)

width = prs.slide_width
height = prs.slide_height
print(f'Slide dimensions: {width}x{height} EMUs ({Emu(width).inches:.2f}x{Emu(height).inches:.2f} inches)')
print(f'Number of slides: {len(prs.slides)}')
print()

# Extract slide master/layout colors
for layout in prs.slide_layouts:
    print(f'Layout: {layout.name}')

print()

img_count = 0
for i, slide in enumerate(prs.slides):
    print(f'=== SLIDE {i+1} ===')
    print(f'Layout: {slide.slide_layout.name}')
    
    for shape in slide.shapes:
        stype = str(shape.shape_type)
        print(f'  Shape: type={stype}, Name="{shape.name}"')
        print(f'    Pos: left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}')
        
        if hasattr(shape, 'text') and shape.text:
            text = shape.text[:300].replace('\n', '\\n')
            print(f'    Text: "{text}"')
            
        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                align = para.alignment
                if align:
                    print(f'      P{pi} align={align}')
                for run in para.runs:
                    font = run.font
                    info = []
                    if font.name: info.append(f'font={font.name}')
                    if font.size: info.append(f'size={font.size.pt}pt')
                    try:
                        if font.color and font.color.rgb: info.append(f'color=#{font.color.rgb}')
                    except: pass
                    if font.bold: info.append('bold')
                    if font.italic: info.append('italic')
                    if info:
                        txt = run.text[:100]
                        print(f'      P{pi} Run: [{txt}] => {", ".join(info)}')
        
        if shape.shape_type == 13:  # Picture
            img_count += 1
            img_data = shape.image.blob
            ext = shape.image.content_type.split('/')[-1]
            if ext == 'jpeg': ext = 'jpg'
            fname = f'slide{i+1}_img{img_count}.{ext}'
            with open(os.path.join(img_dir, fname), 'wb') as f:
                f.write(img_data)
            print(f'    Image saved: {fname} ({shape.image.content_type})')
    print()

# Also extract theme colors
print("=== THEME INFO ===")
theme = prs.slide_masters[0].slide_layouts[0].slide_master.element
# Try to get color scheme
import xml.etree.ElementTree as ET
ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
for clrScheme in theme.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}clrScheme'):
    print(f'Color scheme: {clrScheme.get("name", "unknown")}')
    for child in clrScheme:
        tag = child.tag.split('}')[-1]
        for sub in child:
            val = sub.get('val', sub.get('lastClr', ''))
            print(f'  {tag}: {val}')
